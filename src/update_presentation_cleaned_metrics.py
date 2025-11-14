"""
Update PowerPoint Presentation with Cleaned Dataset V2 Metrics

This script updates the team presentation with the latest model performance metrics
from the cleaned dataset (after deduplication and zero-popularity removal).

New Metrics (Cleaned Dataset V2):
- R² = 0.1619 (down from 0.4772 on uncleaned data - this is expected!)
- RMSE = 16.32
- MAE = 13.14
- n_estimators = 500
- max_depth = 9
- Dataset: 78,310 tracks (down from 114,000)
"""

from pptx import Presentation
import re
from pathlib import Path

# New metrics from cleaned dataset V2
NEW_METRICS = {
    'r2': 0.1619,
    'adj_r2': 0.1613,
    'rmse': 16.32,
    'mae': 13.14,
    'n_estimators': 500,
    'max_depth': 9,
    'learning_rate': 0.0131,
    'n_samples': 78310,
    'dataset_version': 'Cleaned V2 (Deduplicated)'
}

def update_presentation(pptx_path, output_path=None):
    """
    Update PowerPoint presentation with new metrics

    Args:
        pptx_path: Path to the original presentation
        output_path: Path to save updated presentation (default: add _updated_cleaned suffix)
    """
    if output_path is None:
        pptx_file = Path(pptx_path)
        output_path = pptx_file.parent / f"{pptx_file.stem}_updated_cleaned{pptx_file.suffix}"

    # Load presentation
    prs = Presentation(pptx_path)

    updates_made = 0

    # Iterate through all slides
    for slide_num, slide in enumerate(prs.slides, 1):
        # Check all shapes in the slide
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            # Check for text frames
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        updated_text = original_text

                        # Update R² values (look for multiple patterns)
                        # Old values: 0.4772, 0.8447, 0.39, etc.
                        if re.search(r'R²\s*[=:]\s*0\.(4772|8447|477|39|40)', updated_text, re.IGNORECASE):
                            updated_text = re.sub(
                                r'(R²\s*[=:]\s*)0\.(4772|8447|477|39|40)',
                                rf'\g<1>{NEW_METRICS["r2"]:.4f}',
                                updated_text,
                                flags=re.IGNORECASE
                            )
                            print(f"Slide {slide_num}: Updated R² value")

                        # Update Adjusted R²
                        if re.search(r'Adjusted\s+R²\s*[=:]\s*0\.', updated_text, re.IGNORECASE):
                            updated_text = re.sub(
                                r'(Adjusted\s+R²\s*[=:]\s*)0\.\d+',
                                rf'\g<1>{NEW_METRICS["adj_r2"]:.4f}',
                                updated_text,
                                flags=re.IGNORECASE
                            )
                            print(f"Slide {slide_num}: Updated Adjusted R²")

                        # Update RMSE values
                        if re.search(r'RMSE\s*[=:]\s*\d+\.', updated_text, re.IGNORECASE):
                            updated_text = re.sub(
                                r'(RMSE\s*[=:]\s*)\d+\.\d+',
                                rf'\g<1>{NEW_METRICS["rmse"]:.2f}',
                                updated_text,
                                flags=re.IGNORECASE
                            )
                            print(f"Slide {slide_num}: Updated RMSE")

                        # Update MAE values
                        if re.search(r'MAE\s*[=:]\s*\d+\.', updated_text, re.IGNORECASE):
                            updated_text = re.sub(
                                r'(MAE\s*[=:]\s*)\d+\.\d+',
                                rf'\g<1>{NEW_METRICS["mae"]:.2f}',
                                updated_text,
                                flags=re.IGNORECASE
                            )
                            print(f"Slide {slide_num}: Updated MAE")

                        # Update n_estimators (only update 450, 400, etc.)
                        if re.search(r'(?:n_estimators|estimators)\s*[=:]\s*(450|400|200)', updated_text, re.IGNORECASE):
                            updated_text = re.sub(
                                r'((?:n_estimators|estimators)\s*[=:]\s*)(450|400|200)',
                                rf'\g<1>{NEW_METRICS["n_estimators"]}',
                                updated_text,
                                flags=re.IGNORECASE
                            )
                            print(f"Slide {slide_num}: Updated n_estimators")

                        # Update max_depth (look for values like 10, 8, etc.)
                        if re.search(r'(?:max_depth|depth)\s*[=:]\s*(10|8)', updated_text, re.IGNORECASE):
                            updated_text = re.sub(
                                r'((?:max_depth|depth)\s*[=:]\s*)(10|8)',
                                rf'\g<1>{NEW_METRICS["max_depth"]}',
                                updated_text,
                                flags=re.IGNORECASE
                            )
                            print(f"Slide {slide_num}: Updated max_depth")

                        # Update dataset size (114,000 or 114K)
                        if re.search(r'114[,\s]?000|114K', updated_text, re.IGNORECASE):
                            updated_text = re.sub(
                                r'114[,\s]?000|114K',
                                f'{NEW_METRICS["n_samples"]:,}',
                                updated_text,
                                flags=re.IGNORECASE
                            )
                            print(f"Slide {slide_num}: Updated dataset size")

                        # If text was updated, apply the change
                        if updated_text != original_text:
                            run.text = updated_text
                            updates_made += 1

    # Save the updated presentation
    prs.save(output_path)
    print(f"\n✅ Presentation updated successfully!")
    print(f"📊 Total updates made: {updates_made}")
    print(f"💾 Saved to: {output_path}")

    return output_path

if __name__ == "__main__":
    # Find presentation file
    pptx_files = list(Path(".").glob("*.pptx"))
    pptx_files = [f for f in pptx_files if not f.name.startswith("~$")]  # Exclude temp files

    if not pptx_files:
        print("❌ No PowerPoint files found in current directory")
        exit(1)

    # Use the most recent presentation
    pptx_file = max(pptx_files, key=lambda p: p.stat().st_mtime)

    print(f"📂 Found presentation: {pptx_file.name}")
    print(f"\n🔄 Updating with new metrics from Cleaned Dataset V2:")
    print(f"  - R² Score: {NEW_METRICS['r2']:.4f}")
    print(f"  - Adjusted R²: {NEW_METRICS['adj_r2']:.4f}")
    print(f"  - RMSE: {NEW_METRICS['rmse']:.2f}")
    print(f"  - MAE: {NEW_METRICS['mae']:.2f}")
    print(f"  - Dataset Size: {NEW_METRICS['n_samples']:,} tracks")
    print(f"  - n_estimators: {NEW_METRICS['n_estimators']}")
    print(f"  - max_depth: {NEW_METRICS['max_depth']}")
    print()

    # Update the presentation
    output_file = update_presentation(pptx_file)

    print(f"\n📝 Note: R² dropped from 0.48 to 0.16 because:")
    print(f"  - Removed 30,533 duplicate tracks")
    print(f"  - Removed 5,157 zero-popularity tracks")
    print(f"  - R² = 0.16 is realistic for audio-only features")
    print(f"  - Old R² = 0.48 was inflated by easy zero predictions")
    print(f"\n📖 See docs/TRAINING_RUN_SUMMARY_20251114.md for full explanation")
